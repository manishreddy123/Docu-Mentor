import os
import pickle
import hashlib
import pandas as pd
import docx
import pdfplumber
from pptx import Presentation
from agents.chunking_agent import ChunkingAgent
from core.config_manager import ConfigManager

ConfigManager.ensure_directories()


def get_file_hash(filepath):
    with open(filepath, 'rb') as f:
        return hashlib.md5(f.read()).hexdigest()


def parse_documents(filepaths):
    chunking_agent = ChunkingAgent()
    parsed_chunks = []

    for path in filepaths:
        ext = os.path.splitext(path)[-1].lower()
        base_name = os.path.basename(path)

        if ext == ".pdf":
            with pdfplumber.open(path) as pdf:
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        chunks = chunking_agent.chunk(page_text)
                        for chunk in chunks:
                            parsed_chunks.append({
                                "source": f"{base_name} p. {i+1}",
                                "content": chunk,
                                "metadata": {"page": i+1, "filename": base_name}
                            })

        elif ext == ".docx":
            doc = docx.Document(path)
            text = "\n".join([p.text for p in doc.paragraphs])
            chunks = chunking_agent.chunk(text)
            for chunk in chunks:
                parsed_chunks.append({
                    "source": base_name,
                    "content": chunk
                })

        elif ext == ".pptx":
            prs = Presentation(path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            chunks = chunking_agent.chunk(text)
            for chunk in chunks:
                parsed_chunks.append({
                    "source": base_name,
                    "content": chunk
                })

        elif ext == ".csv":
            df = pd.read_csv(path)
            text = df.to_string(index=False)
            chunks = chunking_agent.chunk(text)
            for chunk in chunks:
                parsed_chunks.append({
                    "source": base_name,
                    "content": chunk
                })

        elif ext in [".txt", ".md"]:
            with open(path, "r", encoding="utf-8") as f:
                text = f.read()
            chunks = chunking_agent.chunk(text)
            for chunk in chunks:
                parsed_chunks.append({
                    "source": base_name,
                    "content": chunk
                })
    print(f"✅ Parsed {len(parsed_chunks)} chunks from {len(filepaths)} files")
    return parsed_chunks


def load_or_parse(filepath):
    try:
        file_hash = get_file_hash(filepath)
        cache_path = os.path.join(ConfigManager.CACHE_DIR, f"{file_hash}.pkl")

        if os.path.exists(cache_path):
            try:
                with open(cache_path, 'rb') as f:
                    cached_result = pickle.load(f)
                    print(f"✅ Loaded from cache: {os.path.basename(filepath)}")
                    return cached_result
            except Exception as e:
                print(f"⚠️ Cache read failed for {filepath}: {str(e)}")

        parsed_chunks = parse_documents([filepath])
        
        if not parsed_chunks:
            print(f"⚠️ No parsable content in: {filepath}")
            return [{"source": filepath, "content": ""}]

        try:
            with open(cache_path, 'wb') as f:
                pickle.dump(parsed_chunks, f)
        except Exception as e:
            print(f"⚠️ Cache write failed for {filepath}: {str(e)}")

        return parsed_chunks
    except Exception as e:
        print(f"⚠️ Document loading failed for {filepath}: {str(e)}")
        return [{"source": filepath, "content": ""}]
